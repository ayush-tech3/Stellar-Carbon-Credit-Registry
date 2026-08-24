import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Deep Emerald / Dark Theme
    BG_DARK = RGBColor(10, 25, 20)       # #0a1914
    CARD_BG = RGBColor(18, 38, 30)       # #12261e
    EMERALD = RGBColor(16, 185, 129)     # #10b981
    MINT = RGBColor(52, 211, 153)        # #34d399
    WHITE = RGBColor(245, 247, 250)      # #f5f7fa
    MUTED = RGBColor(156, 163, 175)      # #9ca3af
    ACCENT_CYAN = RGBColor(6, 182, 212)  # #06b6d4

    blank_layout = prs.slide_layouts[6]

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.color.rgb = BG_DARK
        return bg

    def add_header(slide, title_text, category="STELLAR SOROBAN • CARBON REGISTRY"):
        # Category label
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = MINT
        p_cat.font.name = "Arial"

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, title, items, badge=""):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = EMERALD
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        if badge:
            p_b = tf.paragraphs[0]
            p_b.text = badge.upper()
            p_b.font.size = Pt(9)
            p_b.font.bold = True
            p_b.font.color.rgb = ACCENT_CYAN
            p_b.font.name = "Arial"
            p_t = tf.add_paragraph()
        else:
            p_t = tf.paragraphs[0]

        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = MINT
        p_t.font.name = "Arial"
        p_t.space_after = Pt(8)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(11.5)
            p.font.color.rgb = WHITE
            p.font.name = "Arial"
            p.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)

    tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "🌿 CarbonTrack"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = MINT
    p.font.name = "Arial"
    p.space_after = Pt(8)

    p = tf1.add_paragraph()
    p.text = "Decentralized Carbon Credit Registry & Retirement Engine"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.space_after = Pt(16)

    p = tf1.add_paragraph()
    p.text = "Built on Stellar Soroban Smart Contracts • Level 5 Blue Belt Architecture"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Arial"
    p.space_after = Pt(28)

    p = tf1.add_paragraph()
    p.text = "🌐 Live App: carbon-credit-registry.netlify.app   |   📦 GitHub: github.com/ayush-tech3/Stellar-Carbon-Credit-Registry"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.font.name = "Arial"

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_header(slide2, "The Problem: Structural Flaws in Voluntary Carbon Markets")

    add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Double-Counting & Fraud",
             [
                 "Same carbon credits sold to multiple buyers across private silos.",
                 "Zero real-time cryptographic reconciliation between registries.",
                 "Severe greenwashing risk for enterprise ESG audits."
             ], badge="Trust Deficit")

    add_card(slide2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "High Intermediary Costs",
             [
                 "Brokers and middlemen take 15% to 30% of credit value.",
                 "Capital is diverted away from real reforestation & clean energy projects.",
                 "High minimum transaction volumes lock out small participants."
             ], badge="High Friction")

    add_card(slide2, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Slow Reconciliations",
             [
                 "Credit issuance & transfer cycles take days to weeks.",
                 "Manual spreadsheets and PDF certificate verification.",
                 "No immutable, public audit trail for retired credits."
             ], badge="Opaque Systems")

    # -------------------------------------------------------------
    # SLIDE 3: The Solution - CarbonTrack
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_header(slide3, "The Solution: Cryptographic Trust on Stellar Soroban")

    add_card(slide3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "🌱 Verifiable Issuance",
             [
                 "Role-Based Access Control authorizes certified carbon project developers.",
                 "Credits minted with immutable metadata (vintage, methodology, project ID).",
                 "Direct ledger verification prevents duplicate mints."
             ], badge="Issuance Engine")

    add_card(slide3, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "⚡ Sub-Second Transfers",
             [
                 "Direct peer-to-peer and corporate transfers on Stellar.",
                 "Sub-5 second finality with near-zero gas fees (<$0.0001).",
                 "Atomic settlement eliminates counterparty risk."
             ], badge="Liquid Market")

    add_card(slide3, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "🔥 Permanent Retirement",
             [
                 "Atomic burning destroys credits permanently from circulation.",
                 "Mints unique Retirement Certificate (CERT-2026-XXXXX).",
                 "Cryptographically guarantees non-reusability for ESG compliance."
             ], badge="Zero Double-Spend")

    # -------------------------------------------------------------
    # SLIDE 4: System Architecture Overview
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_header(slide4, "End-to-End System Architecture")

    add_card(slide4, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Frontend & UX Tier",
             [
                 "Next.js 15 App Router & React 19.",
                 "Tailwind CSS & Glassmorphic dark theme.",
                 "Zustand State Management for real-time sync.",
                 "Recharts visual analytics & methodology breakdown.",
                 "3D Particle Canvas hero & micro-animations."
             ], badge="Layer 1: Interface")

    add_card(slide4, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Wallet & SDK Tier",
             [
                 "@stellar/stellar-sdk & @stellar/freighter-api.",
                 "1-Click Instant Funded Demo Keypair for frictionless testing.",
                 "Freighter browser wallet integration with explicit permissions.",
                 "Live RPC provider connection to Stellar Soroban Testnet."
             ], badge="Layer 2: Connectivity")

    add_card(slide4, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Smart Contract Tier",
             [
                 "Rust & Soroban SDK on Stellar Testnet.",
                 "Dual-contract architecture with caller verification.",
                 "Contract 1: CarbonCreditRegistry (Mint, Transfer, RBAC).",
                 "Contract 2: RetirementManager (Permanent burn & certs).",
                 "Atomic cross-contract invocations."
             ], badge="Layer 3: Soroban Contracts")

    # -------------------------------------------------------------
    # SLIDE 5: Smart Contract Design & Entrypoints
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_header(slide5, "Soroban Smart Contract Architecture")

    add_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "CarbonCreditRegistry Contract",
             [
                 "initialize(admin) — Sets registry owner & security params.",
                 "add_issuer(issuer) / remove_issuer(issuer) — RBAC governance.",
                 "issue_credits(issuer, project, amount, vintage, methodology) — Mints tokenized credits.",
                 "transfer(from, to, credit_id, amount) — Atomic balances update.",
                 "retire_credits(owner, credit_id, amount, beneficiary) — Calls RetirementManager & burns tokens.",
                 "get_balance(owner, credit_id) & get_total_issued() — Query state."
             ], badge="Core Registry")

    add_card(slide5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "RetirementManager Contract",
             [
                 "record_retirement(credit_id, amount, beneficiary, project) — Records burn record permanently.",
                 "Caller Verification — Restricts invocation strictly to verified Registry Contract.",
                 "Certificate Engine — Generates immutable hash & unique certificate ID (CERT-2026-XXXXX).",
                 "get_retirement_record(cert_id) — Public verifier for auditors.",
                 "get_total_retired() — Global platform aggregate query."
             ], badge="Burn & Certificate Ledger")

    # -------------------------------------------------------------
    # SLIDE 6: Key Product Features & UX Polish
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6)
    add_header(slide6, "Key Product Features & UX Polish")

    add_card(slide6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "1-Click Onboarding",
             [
                 "Funded Demo Wallet mode allows instant testnet interaction.",
                 "Freighter wallet extension with automatic detection.",
                 "No seed phrase hassle for preliminary testing.",
                 "Real-time balance refresh."
             ], badge="Seamless UX")

    add_card(slide6, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Live ESG Analytics",
             [
                 "Interactive Recharts for methodology breakdown (VCS vs Gold Standard).",
                 "Global metrics counters (Total Issued, Retired, Active).",
                 "Real-time live activity feed of testnet events.",
                 "Exportable retirement certificate receipts."
             ], badge="Data & Telemetry")

    add_card(slide6, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Production-Grade Polish",
             [
                 "Mobile responsive drawer & bottom navigation.",
                 "React ErrorBoundary with graceful fallback UI.",
                 "Skeleton loading states for zero layout shift.",
                 "Animated toast alerts for transaction status.",
                 "In-app feedback widget & survey integration."
             ], badge="Enterprise Ready")

    # -------------------------------------------------------------
    # SLIDE 7: User Roles & Workflows
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7)
    add_header(slide7, "Target User Roles & Real-World Workflows")

    add_card(slide7, Inches(0.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "🌱 Project Issuers",
             [
                 "Submit MRV data.",
                 "Mint batch credits directly on-chain.",
                 "Attach vintage & methodology specs.",
                 "Manage inventory in real-time."
             ], badge="Role 1")

    add_card(slide7, Inches(3.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "🏢 Corporate Buyers",
             [
                 "Purchase certified credits peer-to-peer.",
                 "Transfer batches across subsidiaries.",
                 "Sub-second settlement on Stellar."
             ], badge="Role 2")

    add_card(slide7, Inches(6.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "🔥 Net-Zero Officers",
             [
                 "Retire credits permanently.",
                 "Specify beneficiary for ESG reporting.",
                 "Receive immutable CERT-2026 certificate."
             ], badge="Role 3")

    add_card(slide7, Inches(9.8), Inches(1.8), Inches(2.7), Inches(5.0),
             "🔍 ESG Auditors",
             [
                 "Inspect public Stellar ledger.",
                 "Verify retirement hash on Stellar Expert.",
                 "100% guarantee against double counting."
             ], badge="Role 4")

    # -------------------------------------------------------------
    # SLIDE 8: Live Traction & Performance Metrics
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_background(slide8)
    add_header(slide8, "Testnet Traction & Performance Highlights")

    add_card(slide8, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "On-Chain Traction",
             [
                 "1.22M+ tons of CO₂ simulated on Stellar Testnet.",
                 "50+ Onboarded organizations & testnet accounts.",
                 "2 Fully deployed Soroban smart contracts.",
                 "Verified transactions on Stellar Expert Explorer."
             ], badge="Scale")

    add_card(slide8, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Performance & Speed",
             [
                 "Sub-5 second transaction finality.",
                 "<$0.0001 average transaction cost.",
                 "Zero gas price volatility spikes.",
                 "Energy-efficient Stellar Consensus Protocol (SCP)."
             ], badge="Efficiency")

    add_card(slide8, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Code Quality",
             [
                 "100% Soroban Rust unit & integration test pass.",
                 "Comprehensive edge-case coverage.",
                 "Cross-contract authorization safety.",
                 "Clean TypeScript codebase with strict typing."
             ], badge="Reliability")

    # -------------------------------------------------------------
    # SLIDE 9: User Feedback & Iterative Development
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_background(slide9)
    add_header(slide9, "Community Feedback & Product Evolution (Level 4/5)")

    add_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "User Feedback Insights",
             [
                 "• Issuer Feedback: 'Needed an instant preview without connecting browser extension every time.'",
                 "• Auditor Feedback: 'Wanted visual breakdown by methodology (VCS vs Gold Standard).'",
                 "• Enterprise Retiree: 'Freighter wallet permission prompts should be explicit.'",
                 "• Mobile Trader: 'Needed mobile navigation drawer and toast alerts.'"
             ], badge="User Insights")

    add_card(slide9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "Shipped Improvements & Actions",
             [
                 "✅ Built 1-Click Instant Demo Wallet with funded testnet keypair.",
                 "✅ Integrated interactive Recharts methodology analytics chart.",
                 "✅ Added explicit setAllowed()/isAllowed() Freighter handling.",
                 "✅ Implemented mobile responsive drawer, bottom nav, and toast alerts.",
                 "✅ Integrated public Google Feedback Form & live response tracking."
             ], badge="Actions Taken")

    # -------------------------------------------------------------
    # SLIDE 10: Future Roadmap (Level 5 & Beyond)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_background(slide10)
    add_header(slide10, "Future Roadmap & Innovation Horizon")

    add_card(slide10, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Phase 1: Micro-Credits",
             [
                 "Target: Q4 2026",
                 "Upgrade Soroban contracts for micro-fractional credits.",
                 "Enables retail offsetting down to $0.0001 per kg CO₂ for everyday purchases."
             ], badge="Micro-Offsetting")

    add_card(slide10, Inches(4.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Phase 2: Oracle Feeds",
             [
                 "Target: Q1 2027",
                 "Integrate decentralized oracle feeds for Verra & Gold Standard.",
                 "Automatic cross-verification of satellite MRV and real-world registries."
             ], badge="Oracle Layer")

    add_card(slide10, Inches(8.8), Inches(1.8), Inches(3.6), Inches(5.0),
             "Phase 3: DAO Governance",
             [
                 "Target: Q2 2027",
                 "Transition issuer approval (add/remove issuer) to token-weighted DAO voting.",
                 "Fully decentralized community stewardship."
             ], badge="Decentralization")

    # -------------------------------------------------------------
    # SLIDE 11: Summary & Live Links
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_background(slide11)
    add_header(slide11, "Explore CarbonTrack Live")

    add_card(slide11, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.8),
             "Links & Demonstration Resources",
             [
                 "🌐 Live Web App:  https://carbon-credit-registry.netlify.app",
                 "🎥 Video Walkthrough:  https://youtu.be/tyFBRt-QJQs",
                 "📦 GitHub Repository:  https://github.com/ayush-tech3/Stellar-Carbon-Credit-Registry",
                 "📝 User Feedback Form:  https://forms.gle/rF7KsMAaD7SQzQan9",
                 "📊 Public Responses Sheet:  https://docs.google.com/spreadsheets/d/1mKnmxuc9a4YKgHesZjv9jU-2HPNfopE4hsUmRv3Csp4",
                 "🔍 Stellar Expert Explorer:  https://stellar.expert/explorer/testnet"
             ], badge="Live Deployment")

    output_path = "CarbonTrack_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_deck()
